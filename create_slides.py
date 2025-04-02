import sys
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def create_professional_slides(markdown_content, output_file, template_file=None):
    """
    Converts markdown content to professional PowerPoint slides
    
    Args:
        markdown_content (str): Markdown text with slide separators (---)
        output_file (str): Filename for output PowerPoint
        template_file (str, optional): PowerPoint template to use
    """
    # Use template if provided, otherwise create new presentation
    if template_file and os.path.exists(template_file):
        prs = Presentation(template_file)
    else:
        prs = Presentation()
    
    # Set consistent font sizes and styles
    TITLE_FONT_SIZE = Pt(32)
    SUBTITLE_FONT_SIZE = Pt(20)
    CONTENT_FONT_SIZE = Pt(18)
    BULLET_FONT_SIZE = Pt(18)
    
    # Professional color scheme
    TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
    TEXT_COLOR = RGBColor(0, 0, 0)      # Black
    ACCENT_COLOR = RGBColor(0, 102, 204) # Medium blue
    
    sections = re.split(r'\n---\n', markdown_content)
    
    for section in sections:
        if not section.strip():
            continue
            
        lines = section.strip().split('\n')
        
        # Find title (lines with # or ##)
        title_line = next((line for line in lines if re.match(r'^#+\s+', line)), None)
        
        # Skip section if no title found
        if not title_line:
            continue
            
        # Determine slide layout based on content
        if title_line.startswith('# '):  # Main title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
            title_text = title_line.replace('# ', '')
            subtitle_line = lines[1] if len(lines) > 1 and not lines[1].startswith('#') else ""
            
            # Set title with formatting
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.alignment = PP_ALIGN.CENTER
            for run in title_para.runs:
                run.font.size = TITLE_FONT_SIZE
                run.font.color.rgb = TITLE_COLOR
                run.font.bold = True
            
            # Set subtitle with formatting if it exists
            if len(slide.placeholders) > 1 and subtitle_line:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle_line
                subtitle_para = subtitle_shape.text_frame.paragraphs[0]
                subtitle_para.alignment = PP_ALIGN.CENTER
                for run in subtitle_para.runs:
                    run.font.size = SUBTITLE_FONT_SIZE
                    run.font.color.rgb = TEXT_COLOR
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
            title_text = title_line.replace('## ', '')
            
            # Set title with formatting
            title_shape = slide.shapes.title
            title_shape.text = title_text
            title_para = title_shape.text_frame.paragraphs[0]
            for run in title_para.runs:
                run.font.size = TITLE_FONT_SIZE
                run.font.color.rgb = TITLE_COLOR
                run.font.bold = True
            
            # Collect content (excluding the title)
            content_lines = [line for line in lines if line != title_line]
            
            # Process content for bullets and formatting
            if content_lines and len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                
                p = None
                for line in content_lines:
                    # Skip empty lines
                    if not line.strip():
                        continue
                    
                    # Bullet points with proper indentation
                    if line.strip().startswith('- '):
                        p = text_frame.add_paragraph()
                        p.text = line.strip()[2:]
                        p.level = 0
                        for run in p.runs:
                            run.font.size = BULLET_FONT_SIZE
                            run.font.color.rgb = TEXT_COLOR
                    
                    # Sub-bullets with proper indentation
                    elif line.strip().startswith('  - '):
                        p = text_frame.add_paragraph()
                        p.text = line.strip()[4:]
                        p.level = 1
                        for run in p.runs:
                            run.font.size = BULLET_FONT_SIZE
                            run.font.color.rgb = TEXT_COLOR
                    
                    # Regular text
                    else:
                        p = text_frame.add_paragraph()
                        p.text = line
                        for run in p.runs:
                            run.font.size = CONTENT_FONT_SIZE
                            run.font.color.rgb = TEXT_COLOR
    
    prs.save(output_file)
    print(f"Successfully created professional slides: {output_file}")


def process_all_markdown_files(input_dir=None, output_dir=None, template_file=None):
    """
    Process all markdown files in the input directory and convert them to PowerPoint
    
    Args:
        input_dir (str, optional): Directory containing markdown files (default: current dir)
        output_dir (str, optional): Directory for output files (default: current dir)
        template_file (str, optional): PowerPoint template to use
    """
    if input_dir is None:
        input_dir = os.getcwd()
    
    if output_dir is None:
        output_dir = os.getcwd()
    
    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Find all markdown files
    markdown_files = glob.glob(os.path.join(input_dir, "*.md"))
    
    if not markdown_files:
        print(f"No markdown files found in {input_dir}")
        return
    
    for md_file in markdown_files:
        base_name = os.path.basename(md_file)
        base_name_no_ext = os.path.splitext(base_name)[0]
        output_file = os.path.join(output_dir, f"{base_name_no_ext}.pptx")
        
        # Skip README.md and workflow documents
        if base_name.lower() in ["readme.md", "workflow.md", "workflow-system.md"]:
            continue
        
        try:
            with open(md_file, 'r') as f:
                markdown_content = f.read()
            
            create_professional_slides(markdown_content, output_file, template_file)
        except Exception as e:
            print(f"Error processing {md_file}: {str(e)}")


# Handle command line arguments
if __name__ == "__main__":
    if len(sys.argv) == 1:
        # No arguments, process all markdown files in current directory
        process_all_markdown_files()
    elif len(sys.argv) == 2 and (sys.argv[1] == "-h" or sys.argv[1] == "--help"):
        print("Usage:")
        print("  python create_slides.py                       # Process all markdown files in current directory")
        print("  python create_slides.py input.md output.pptx  # Process a single markdown file")
        print("  python create_slides.py input_dir output_dir  # Process all markdown files in input_dir")
        print("  python create_slides.py input_dir output_dir template.pptx  # Use a PowerPoint template")
    elif len(sys.argv) == 3:
        # Two arguments, could be:
        # 1. Single file conversion: input.md output.pptx
        # 2. Batch processing: input_dir output_dir
        if os.path.isfile(sys.argv[1]):
            # Single file conversion
            input_file = sys.argv[1]
            output_file = sys.argv[2]
            
            try:
                with open(input_file, 'r') as f:
                    markdown_content = f.read()
                
                create_professional_slides(markdown_content, output_file)
            except Exception as e:
                print(f"Error: {str(e)}")
                sys.exit(1)
        elif os.path.isdir(sys.argv[1]):
            # Batch processing
            input_dir = sys.argv[1]
            output_dir = sys.argv[2]
            
            process_all_markdown_files(input_dir, output_dir)
        else:
            print(f"Error: {sys.argv[1]} is neither a file nor a directory.")
            sys.exit(1)
    elif len(sys.argv) == 4:
        # Three arguments: input_dir output_dir template.pptx
        input_dir = sys.argv[1]
        output_dir = sys.argv[2]
        template_file = sys.argv[3]
        
        if not os.path.isdir(input_dir):
            print(f"Error: {input_dir} is not a directory.")
            sys.exit(1)
        
        if not os.path.isfile(template_file):
            print(f"Error: Template file {template_file} not found.")
            sys.exit(1)
        
        process_all_markdown_files(input_dir, output_dir, template_file)
    else:
        # Default filenames if not provided
        input_file = "slides.md"
        output_file = "slides.pptx"
        
        try:
            # Read markdown content
            with open(input_file, 'r') as f:
                markdown_content = f.read()
            
            # Create slides
            create_professional_slides(markdown_content, output_file)
        except Exception as e:
            print(f"Error: {str(e)}")
            sys.exit(1)